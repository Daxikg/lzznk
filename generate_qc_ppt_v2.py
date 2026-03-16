# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml

# 创建PPT
prs = Presentation()
prs.slide_width = Inches(13.333)  # 16:9比例
prs.slide_height = Inches(7.5)

def set_shape_fill(shape, r, g, b):
    """设置形状背景色"""
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(r, g, b)

def add_title_slide(prs, title, subtitle='', bottom_text=''):
    """添加标题幻灯片"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局

    # 添加深色背景条
    shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(2.2))
    set_shape_fill(shape, 0, 51, 102)  # 深蓝色

    # 添加主标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.6), Inches(12.333), Inches(1.2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.name = '黑体'
    p.font.color.rgb = RGBColor(255, 255, 255)  # 白色
    p.alignment = PP_ALIGN.CENTER

    # 添加副标题
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.0), Inches(12.333), Inches(1.5))
        tf2 = sub_box.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(28)
        p2.font.name = '宋体'
        p2.alignment = PP_ALIGN.CENTER

    # 添加底部文字
    if bottom_text:
        bottom_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12.333), Inches(0.5))
        tf3 = bottom_box.text_frame
        p3 = tf3.paragraphs[0]
        p3.text = bottom_text
        p3.font.size = Pt(18)
        p3.font.name = '宋体'
        p3.alignment = PP_ALIGN.CENTER

    return slide

def add_section_slide(prs, section_num, section_title):
    """添加章节标题幻灯片"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 添加左侧深色条
    shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.3), Inches(7.5))
    set_shape_fill(shape, 0, 51, 102)

    # 添加章节编号
    num_box = slide.shapes.add_textbox(Inches(1), Inches(2.8), Inches(2), Inches(1.5))
    tf = num_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"0{section_num}" if section_num < 10 else str(section_num)
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.name = 'Arial'
    p.font.color.rgb = RGBColor(0, 51, 102)

    # 添加章节标题
    title_box = slide.shapes.add_textbox(Inches(3.5), Inches(3.0), Inches(9), Inches(1.2))
    tf2 = title_box.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = section_title
    p2.font.size = Pt(40)
    p2.font.bold = True
    p2.font.name = '黑体'

    return slide

def add_content_slide(prs, title, content_items, hint_image=''):
    """添加内容幻灯片"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 添加顶部深色条
    shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(1.0))
    set_shape_fill(shape, 0, 51, 102)

    # 添加标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.333), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.name = '黑体'
    p.font.color.rgb = RGBColor(255, 255, 255)

    # 添加内容
    content_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.3), Inches(12.133), Inches(5.5))
    tf2 = content_box.text_frame
    tf2.word_wrap = True

    for i, item in enumerate(content_items):
        if i == 0:
            p = tf2.paragraphs[0]
        else:
            p = tf2.add_paragraph()
        p.text = item
        p.font.size = Pt(18)
        p.font.name = '宋体'
        p.space_after = Pt(8)

    # 添加图片提示
    if hint_image:
        hint_box = slide.shapes.add_textbox(Inches(0.6), Inches(6.5), Inches(12.133), Inches(0.5))
        tf3 = hint_box.text_frame
        p3 = tf3.paragraphs[0]
        p3.text = f"【建议插入图片：{hint_image}】"
        p3.font.size = Pt(14)
        p3.font.name = '宋体'
        p3.font.color.rgb = RGBColor(255, 0, 0)

    return slide

def add_table_slide(prs, title, headers, data, hint_image=''):
    """添加表格幻灯片"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 添加顶部深色条
    shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(1.0))
    set_shape_fill(shape, 0, 51, 102)

    # 添加标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.333), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.name = '黑体'
    p.font.color.rgb = RGBColor(255, 255, 255)

    # 添加表格
    rows = len(data) + 1
    cols = len(headers)
    table = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.2), Inches(12.333), Inches(5.0)).table

    # 设置表头
    for j, header in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = header
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.bold = True
            paragraph.font.size = Pt(14)
            paragraph.font.name = '黑体'
            paragraph.alignment = PP_ALIGN.CENTER

    # 填充数据
    for i, row_data in enumerate(data):
        for j, cell_text in enumerate(row_data):
            cell = table.cell(i+1, j)
            cell.text = str(cell_text)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(12)
                paragraph.font.name = '宋体'

    # 添加图片提示
    if hint_image:
        hint_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12.333), Inches(0.5))
        tf3 = hint_box.text_frame
        p3 = tf3.paragraphs[0]
        p3.text = f"【建议插入图片：{hint_image}】"
        p3.font.size = Pt(14)
        p3.font.name = '宋体'
        p3.font.color.rgb = RGBColor(255, 0, 0)

    return slide

def add_image_placeholder_slide(prs, title, hint_text):
    """添加图片占位幻灯片"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 添加顶部深色条
    shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(1.0))
    set_shape_fill(shape, 0, 51, 102)

    # 添加标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.333), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.name = '黑体'
    p.font.color.rgb = RGBColor(255, 255, 255)

    # 添加图片占位框
    placeholder = slide.shapes.add_shape(1, Inches(1.5), Inches(1.5), Inches(10.333), Inches(5.0))
    fill = placeholder.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(240, 240, 240)

    # 添加提示文字
    hint_box = slide.shapes.add_textbox(Inches(1.5), Inches(3.5), Inches(10.333), Inches(1))
    tf2 = hint_box.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = f"【插入图片：{hint_text}】"
    p2.font.size = Pt(24)
    p2.font.name = '宋体'
    p2.font.color.rgb = RGBColor(128, 128, 128)
    p2.alignment = PP_ALIGN.CENTER

    return slide

# ==================== 幻灯片1：封面 ====================
add_title_slide(prs, '贵阳南检修信息管理平台', '贵阳南车辆段\n检修车间QC小组', '2025年01月—2025年12月')

# ==================== 幻灯片2：QC小组概况 ====================
add_content_slide(prs, 'QC小组概况', [
    '• 小组名称：检修车间QC小组',
    '• 成立日期：2025年01月',
    '• 活动时间：2025.01-2025.12',
    '• 活动次数：24次',
    '• 课题名称：贵阳南检修信息管理平台',
    '• 课题类型：创新型',
    '• 小组成员：10人',
    '  - 组长1人（车间主任）',
    '  - 副组长2人（党总支书记、车间副主任）',
    '  - 技术指导员2人',
    '  - 组员5人',
    '• 平均受QC教育时间：24学时'
])

# ==================== 幻灯片3：小组成员名单 ====================
add_table_slide(prs, '小组成员名单',
    ['序号', '姓名', '职务', '职称', '组内分工'],
    [
        ['1', '（组长）', '车间主任', '工程师', '总体方案指导'],
        ['2', '（副组长）', '车间党总支书记', '政工师', '方案指导'],
        ['3', '（副组长）', '车间副主任', '高级工程师', '方案策划'],
        ['4', '（技术指导员）', '技术员', '助理工程师', '计划制定'],
        ['5', '（技术指导员）', '技术员', '助理工程师', '软件开发'],
        ['6-10', '（组员）', '班组骨干', '助理工程师/技师', '方案实施'],
    ],
    '小组成员合影照片'
)

# ==================== 幻灯片4：目录 ====================
slide4 = prs.slides.add_slide(prs.slide_layouts[6])

# 添加顶部深色条
shape = slide4.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(1.0))
set_shape_fill(shape, 0, 51, 102)

title_box = slide4.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.333), Inches(0.7))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.text = '目  录'
p.font.size = Pt(28)
p.font.bold = True
p.font.name = '黑体'
p.font.color.rgb = RGBColor(255, 255, 255)

# 目录项
toc_items = [
    ('01', '选题', '选题背景、广泛借鉴、设定目标'),
    ('02', '提出方案并确定最佳方案', '总体设计、功能模块、可行性论证'),
    ('03', '设定目标及目标可行性论证', '课题目标、可行性分析'),
    ('04', '制定对策', '需求调研、模块划分'),
    ('05', '对策实施', '功能开发、系统测试'),
    ('06', '效果检查', '功能验证、效益分析'),
    ('07', '总结和下一步打算', '成果总结、推广应用'),
]

y_pos = 1.3
for num, title, desc in toc_items:
    # 编号
    num_box = slide4.shapes.add_textbox(Inches(1), Inches(y_pos), Inches(1), Inches(0.6))
    tf = num_box.text_frame
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.name = 'Arial'
    p.font.color.rgb = RGBColor(0, 51, 102)

    # 标题
    title_box2 = slide4.shapes.add_textbox(Inches(2.2), Inches(y_pos), Inches(4), Inches(0.6))
    tf2 = title_box2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = title
    p2.font.size = Pt(20)
    p2.font.bold = True
    p2.font.name = '黑体'

    # 描述
    desc_box = slide4.shapes.add_textbox(Inches(6.5), Inches(y_pos), Inches(6), Inches(0.6))
    tf3 = desc_box.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = desc
    p3.font.size = Pt(16)
    p3.font.name = '宋体'
    p3.font.color.rgb = RGBColor(100, 100, 100)

    y_pos += 0.8

# ==================== 幻灯片5：前言说明 ====================
add_section_slide(prs, 1, '前言说明')

# ==================== 幻灯片6：前言内容 ====================
add_content_slide(prs, '前  言', [
    '随着铁路货车检修工作不断发展，车间管理工作日益繁重，各类数据分散存储、',
    '查阅不便、统计困难等问题日益凸显。传统纸质记录和分散的信息系统已无法满足',
    '现代化检修车间高效管理的需求。',
    '',
    '在集团公司党委、集团公司坚强领导下，贵阳南车辆段秉持"匠心至善、铿锵致远"',
    '的南辆精神，围绕"紧盯车间、班组、岗位发力，解决标准落实最后一公里的问题"',
    '的工作目标，积极探索，历时一年，自主研发了满足智能化、一体化新型管理要求',
    '的"贵阳南检修信息管理平台"。',
    '',
    '平台以共享互通为基本要义，解决工作中缺乏互联互通、数据整合、共享交换等问题，',
    '打破"孤岛"效应，实现管理过程、结果透明可视。'
])

# ==================== 幻灯片7：选题背景 ====================
add_section_slide(prs, 2, '选题背景')

# ==================== 幻灯片8：选题背景详细 ====================
add_content_slide(prs, '一、选题背景', [
    '随着铁路货车检修业务的不断发展，检修车间日常管理工作面临着诸多挑战：',
    '',
    '（1）数据分散存储',
    '    各班组、各业务模块的数据分散存储在不同系统和纸质档案中，无法形成完整',
    '    的数据分析，查阅极为不便。',
    '',
    '（2）信息孤岛问题',
    '    车间现有多个信息系统，但各系统之间相互独立，关键业务数据无法共享传递。',
    '',
    '（3）管理效率低下',
    '    班组日志、健康管理、党建管理等日常管理工作依赖纸质记录，统计汇总耗时耗力。'
])

# ==================== 幻灯片9：选题背景续 ====================
add_content_slide(prs, '一、选题背景（续）', [
    '（4）生产协调困难',
    '    生产计划、轮轴选配、物料管理等生产相关信息的传递和协调主要依靠人工方式，',
    '    效率低且容易出错。',
    '',
    '（5）一线职工负担重',
    '    重复填报、多头报送等问题增加了一线职工的工作负担，影响了检修工作的正常开展。',
    '',
    '',
    '以上问题制约了车间管理效率的提升，迫切需要一套综合信息管理平台来解决这些问题。'
])

# ==================== 幻灯片10：广泛借鉴 ====================
add_content_slide(prs, '广泛借鉴，启发创新灵感', [
    '为解决上述问题，QC小组成员广泛调研了国内外各类信息管理系统的设计理念和实现方式：',
    '',
    '（1）现代Web开发技术',
    '    采用Python语言和Django框架，具有开发效率高、安全性好、可扩展性强等优点。',
    '',
    '（2）微服务架构思想',
    '    将系统划分为多个独立模块，降低耦合度，便于维护和扩展。',
    '',
    '（3）用户体验设计',
    '    参考优秀管理平台的界面设计，力求简洁美观、操作便捷。',
    '',
    '（4）数据整合技术',
    '    建立统一数据中心，实现多源数据的集中存储和共享利用。'
])

# ==================== 幻灯片11：设定目标 ====================
add_table_slide(prs, '设定目标',
    ['目标类型', '具体目标'],
    [
        ['功能目标', '开发一套涵盖车间管理、生产系统、规章制度三大类的综合信息管理平台'],
        ['技术目标', '采用Python+Django+MySQL技术栈，实现B/S架构，支持多终端访问'],
        ['效率目标', '将数据查询效率提升80%以上，减少重复填报工作50%以上'],
        ['推广目标', '系统一次部署即可实现车间管理人员、各班组录入人员二级管理模式'],
    ]
)

# ==================== 幻灯片12：提出方案 ====================
add_section_slide(prs, 3, '提出方案并确定最佳方案')

# ==================== 幻灯片13：总体方案设计 ====================
add_content_slide(prs, '总体方案设计', [
    '方案名称：贵阳南检修信息管理平台',
    '',
    '技术架构：',
    '    • 采用B/S架构',
    '    • 后端使用Python Django框架',
    '    • 数据库使用MySQL',
    '    • 前端采用HTML5+CSS3+JavaScript技术栈',
    '',
    '部署方式：',
    '    • 部署在车间服务器',
    '    • 通过局域网访问',
    '    • 支持PC端和移动端（手持机）访问'
], '系统架构图')

# ==================== 幻灯片14：功能模块设计 ====================
add_table_slide(prs, '功能模块设计',
    ['模块类别', '模块名称', '功能说明'],
    [
        ['车间管理类', '车间构架/班组日志/健康管理/消防管理/党建管理', '日常管理功能'],
        ['生产系统类', '物料管理/轮轴选配/设备管理/预检信息/HMIS系统等', '生产辅助功能'],
        ['规章制度类', '设备操作规程/规章查询/作业指导书', '制度文档管理'],
    ],
    '功能模块架构图'
)

# ==================== 幻灯片15：详细功能模块 ====================
add_content_slide(prs, '功能模块详情（共20个模块）', [
    '车间管理类（5个模块）：',
    '    车间构架 | 班组日志 | 健康管理 | 消防管理 | 党建管理',
    '',
    '生产系统类（12个模块）：',
    '    物料管理系统 | 轮轴智能选配系统 | 设备管理信息系统 | 检修车预检信息',
    '    综合视频调阅 | 成都局管理信息系统 | 电子一车一档 | 检修HMIS系统',
    '    轮轴HMIS系统 | 站修HMIS系统 | 铁路机辆物料管理 | HMIS统计分析子系统',
    '',
    '规章制度类（3个模块）：',
    '    设备操作规程 | 规章查询 | 作业指导书'
])

# ==================== 幻灯片16：可行性论证 ====================
add_table_slide(prs, '方案可行性论证',
    ['论证项目', '论证内容', '结论'],
    [
        ['需求分析', '对近段时间以来收集的现场职工诉求进行汇总分析，需求明确', '可行'],
        ['技术能力', '开发人员具有丰富的软件开发经验，技术储备充足', '可行'],
        ['资源保障', '车间提供服务器、网络等基础设施支持', '可行'],
        ['时间安排', '项目周期12个月，时间充足', '可行'],
    ]
)

# ==================== 幻灯片17：制定对策 ====================
add_section_slide(prs, 4, '制定对策')

# ==================== 幻灯片18：对策表 ====================
add_table_slide(prs, '制定对策',
    ['序号', '对策', '目标', '完成时间'],
    [
        ['1', '需求调研', '明确功能需求', '2025.01'],
        ['2', '系统设计', '完成系统架构设计', '2025.02'],
        ['3', '前端开发', '完成前端页面设计', '2025.04'],
        ['4', '后端开发', '完成核心功能开发', '2025.06'],
        ['5', '系统集成', '完成系统集成测试', '2025.07'],
        ['6', '培训推广', '完成用户培训', '2025.08'],
        ['7', '效果检查', '验证系统效果', '2025.09'],
        ['8', '推广应用', '扩大应用范围', '2025.10'],
    ]
)

# ==================== 幻灯片19：对策实施 ====================
add_section_slide(prs, 5, '对策实施')

# ==================== 幻灯片20：实施一 ====================
add_content_slide(prs, '实施一：需求调研与分析', [
    '按前期对"检修信息管理平台"的讨论结果，进行了班组需求调研，并进行了汇总工作，',
    '形成了最终待开发迭代需求表。',
    '',
    '主要需求包括：',
    '',
    '（1）班组日志模块：支持班组人员信息管理、活动记录填写、照片上传等功能',
    '',
    '（2）健康管理模块：支持健康风险人员动态监测、检测提醒、历史记录查询',
    '',
    '（3）党建管理模块：支持党支部信息展示、党员积分管理、创岗建区评定',
    '',
    '（4）轮轴选配模块：支持轮对数据录入、智能配轮计算、结果展示',
    '',
    '（5）物料管理模块：支持物料申请、核销、库存管理等功能'
], '需求调研现场照片或需求汇总表截图')

# ==================== 幻灯片21：实施二 ====================
add_content_slide(prs, '实施二：系统架构设计', [
    '技术选型：',
    '',
    '    • 后端框架：Django 5.0（Python 3.x）',
    '',
    '    • 数据库：MySQL 8.0',
    '',
    '    • 前端技术：HTML5 + CSS3 + JavaScript + Bootstrap',
    '',
    '    • 部署环境：Windows Server + IIS',
    '',
    '系统架构：采用B/S架构，支持多终端访问，包括PC端和手持机端。'
], '系统架构图')

# ==================== 幻灯片22：实施三-车间管理类 ====================
add_content_slide(prs, '实施三：核心功能开发——车间管理类', [
    '（1）班组日志模块',
    '    实现了班组人员信息查询与维护、活动记录填写、照片上传等功能。',
    '',
    '（2）健康管理模块',
    '    实现了健康风险人员动态监测功能，支持每月/每周/每日检测提醒。',
    '',
    '（3）党建管理模块',
    '    实现了检修车间党支部介绍、支委概况、组织结构、创岗建区展示等功能。',
    '',
    '（4）消防管理模块',
    '    实现了设备二维码扫码点检、点检记录实时上传、负责人列表维护等功能。'
], '班组日志/健康管理/党建管理模块界面截图')

# ==================== 幻灯片23：实施三-生产系统类 ====================
add_content_slide(prs, '实施三：核心功能开发——生产系统类', [
    '（1）轮轴智能选配系统',
    '    实现了架车班数据录入、轮轴班数据录入、配轮数据自动计算、轮对选配展示等',
    '    功能，支持PC端和手持机端操作。',
    '',
    '（2）物料管理系统',
    '    实现了物料申请、核销、库存管理等功能，支持材料员、配送员、管理员多角色操作。',
    '',
    '（3）检修车预检信息系统',
    '    实现了车辆预检、预修计划、入线计划、车辆竣工等功能，支持可视化拖拽操作。'
], '轮轴选配/物料管理系统界面截图')

# ==================== 幻灯片24：实施三-基础功能 ====================
add_content_slide(prs, '实施三：核心功能开发——基础功能', [
    '（1）用户权限管理',
    '    实现了基于角色的权限控制，支持管理员、班组账号等多种角色，不同角色具有',
    '    不同的操作权限。',
    '',
    '（2）重要事务通知',
    '    实现了事务通知发布、附件上传、已读未读状态显示等功能。',
    '',
    '（3）文件传阅',
    '    实现了记名式文件传阅功能，支持拍照确认。'
], '系统登录/重要事务通知界面截图')

# ==================== 幻灯片25：实施四 ====================
add_content_slide(prs, '实施四：系统集成与测试', [
    '完成了各模块的集成测试，重点测试了以下内容：',
    '',
    '（1）功能测试',
    '    验证各功能模块是否符合需求文档要求。',
    '',
    '（2）性能测试',
    '    测试系统响应速度、并发处理能力。',
    '',
    '（3）兼容性测试',
    '    测试PC端和移动端的兼容性。',
    '',
    '（4）安全测试',
    '    测试用户权限控制、数据安全等方面。',
    '',
    '针对测试发现的问题，及时进行了修复和优化。'
], '系统测试现场照片')

# ==================== 幻灯片26：效果检查 ====================
add_section_slide(prs, 6, '效果检查')

# ==================== 幻灯片27：功能验证 ====================
add_table_slide(prs, '功能验证',
    ['序号', '功能模块', '验证结果'],
    [
        ['1', '平台内提供的链接可用性', '✓ 正常'],
        ['2', '作业指导书模块应用', '✓ 正常'],
        ['3', '班组日志模块应用', '✓ 正常'],
        ['4', '健康管理模块应用', '✓ 正常'],
        ['5', '党建管理模块应用', '✓ 正常'],
        ['6', '生产进度展示', '✓ 正常'],
        ['7', '重要事务通知', '✓ 正常'],
        ['8', '轮轴智能选配', '✓ 正常'],
        ['9', '物料管理系统', '✓ 正常'],
        ['10', '用户权限控制', '✓ 正常'],
        ['11', '移动端适配', '✓ 正常'],
    ]
)

# ==================== 幻灯片28：目标达成情况 ====================
add_table_slide(prs, '目标达成情况',
    ['目标项目', '目标值', '达成情况'],
    [
        ['功能模块数量', '≥15个', '已开发20个模块，超额完成'],
        ['数据查询效率提升', '≥80%', '达到预期，查询效率显著提升'],
        ['减少重复填报', '≥50%', '达到预期，大幅减少重复工作'],
        ['系统用户覆盖', '车间全员', '已覆盖所有班组和管理人员'],
        ['多终端支持', 'PC端+移动端', '已实现，支持手持机操作'],
    ]
)

# ==================== 幻灯片29：效益分析 ====================
add_content_slide(prs, '效益分析', [
    '经济效益：',
    '    • 降低管理成本：通过信息化手段，减少了纸质资料的使用和存储成本',
    '    • 提高工作效率：数据查询效率提升80%以上，大幅减少人工统计时间',
    '    • 减少人力投入：班组日志、健康管理等工作的信息化，减少了重复填报工作',
    '',
    '社会效益：',
    '    • 提升管理规范化水平：统一的平台和标准化的流程',
    '    • 改善职工体验：便捷的信息查询和填报方式',
    '    • 促进信息共享：打破了信息孤岛，实现了数据的有效共享和利用',
    '',
    '技术效益：',
    '    • 积累软件开发经验：通过项目开发，积累了Web应用开发的经验和技术',
    '    • 形成可推广成果：系统的成功开发，为其他车间的信息化建设提供了参考'
])

# ==================== 幻灯片30：系统展示 ====================
add_image_placeholder_slide(prs, '系统主界面展示', '系统主界面截图')

# ==================== 幻灯片31：总结 ====================
add_section_slide(prs, 7, '总结和下一步打算')

# ==================== 幻灯片32：总结内容 ====================
add_content_slide(prs, '总  结', [
    '通过本次QC活动，成功开发了"贵阳南检修信息管理平台"，实现了以下成果：',
    '',
    '（1）建立了集中存储、标准统一的数据中心，实现了车间管理数据的集中存储和',
    '    共享利用；通过信息化管理，减少了纸质材料消耗，降低了管理成本。',
    '',
    '（2）开发和集成了20个功能模块，涵盖车间管理、生产系统、规章制度三大类，',
    '    满足了车间日常管理和生产的需要。',
    '',
    '（3）实现了PC端和移动端的多终端访问，方便了现场作业人员的使用。',
    '',
    '（4）显著提升了工作效率，减少了重复填报工作，降低了职工工作负担。',
    '',
    '（5）通过活动开展，提高了小组成员的软件开发能力和QC活动开展水平。'
])

# ==================== 幻灯片33：下一步打算 ====================
add_content_slide(prs, '下一步打算', [
    '（1）持续优化',
    '    继续跟进用户反馈内容、意见，对"管理平台"各模块进行优化，确保符合车间使用要求。',
    '',
    '（2）推广应用',
    '    在段级不同车间进行推广、宣传、试用，进一步扩大"管理平台"的影响范围。',
    '',
    '（3）功能扩展',
    '    根据用户需求，持续增加新的功能模块，如数据分析、智能预警等。',
    '',
    '（4）系统升级',
    '    考虑采用更先进的技术架构，提升系统的性能和安全性。',
    '',
    '（5）标准建设',
    '    将成功经验形成技术标准和管理制度，为其他单位提供参考。'
])

# ==================== 幻灯片34：结束页 ====================
add_title_slide(prs, '感谢聆听！', '欢迎各位领导、专家批评指正', '')

# 保存PPT
output_path = r"D:\桌面\QC\贵阳南检修信息管理平台-QC成果报告.pptx"
prs.save(output_path)
print(f"PPT已保存到: {output_path}")